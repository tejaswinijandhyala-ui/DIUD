import csv
import io
import json
import os
import re
import traceback
import uuid
from datetime import date, datetime
from typing import Dict, List, Literal, Optional

import httpx
import anthropic
from fastapi import FastAPI, Header, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate,
    PageBreak, Paragraph, Spacer, Table, TableStyle,
)

# =============================================================================
# Load ENV
# =============================================================================
load_dotenv()

# =============================================================================
# FastAPI App
# =============================================================================
app = FastAPI(title="DIUD", description="Decision Intelligence Using Data", version="4.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =============================================================================
# Claude client
# =============================================================================
_ai_client    = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
_CLAUDE_MODEL = "claude-sonnet-4-5"

# =============================================================================
# SERVER-SIDE SESSION STORE
# Keeps the last query result per session so export endpoints can always
# access the full raw dataset — no JSON embedding in messages, no truncation.
# Key  : session_id (UUID string, generated once per browser tab)
# Value: QueryResult dict with full rows + metadata
# =============================================================================

class QueryResult:
    """Holds the full raw result of the most recent ClickHouse query in a session."""
    def __init__(self, sql: str, columns: List[str], rows: List[dict],
                 total_rows: int, captured_at: str, filters_applied: str = ""):
        self.sql            = sql
        self.columns        = columns
        self.rows           = rows          # ALL rows, no cap
        self.total_rows     = total_rows
        self.captured_at    = captured_at
        self.filters_applied = filters_applied

_SESSION_STORE: Dict[str, QueryResult] = {}   # session_id → QueryResult

def _store_result(session_id: str, result: QueryResult):
    _SESSION_STORE[session_id] = result

def _get_result(session_id: str) -> Optional[QueryResult]:
    return _SESSION_STORE.get(session_id)


# =============================================================================
# ClickHouse HTTP proxy — base helpers
# =============================================================================
def _base_url() -> str:
    return (os.getenv("CLICKHOUSE_API_URL") or "").rstrip("/")

def _token() -> str:
    return os.getenv("CLICKHOUSE_API_TOKEN") or ""

def _auth_headers() -> dict:
    return {
        "Authorization": f"Bearer {_token()}",
        "Content-Type":  "application/json",
    }

FORBIDDEN_KEYWORDS = ["INSERT", "UPDATE", "DELETE", "DROP", "ALTER", "TRUNCATE", "CREATE"]

# =============================================================================
# Schema discovery
# =============================================================================
_LIVE_SCHEMA: dict = {}
_SCHEMA_BLOCK: str = "Schema not yet loaded."


def _proxy_get(path: str) -> dict | list | None:
    base_url = _base_url()
    token    = _token()
    if not base_url or not token:
        return None
    try:
        r = httpx.get(f"{base_url}{path}", headers=_auth_headers(), timeout=20)
        if r.status_code == 200:
            return r.json()
        print(f"   ⚠️  GET {path} → HTTP {r.status_code}: {r.text[:200]}")
        return None
    except Exception as e:
        print(f"   ⚠️  GET {path} → {e}")
        return None


def discover_schema() -> str:
    global _LIVE_SCHEMA, _SCHEMA_BLOCK

    print("🔎 Discovering schema from ClickHouse proxy…")
    databases_raw = _proxy_get("/databases")
    if not databases_raw:
        msg = "⚠️  Could not fetch databases — check CLICKHOUSE_API_URL / CLICKHOUSE_API_TOKEN."
        print(msg); _SCHEMA_BLOCK = msg; return msg

    if isinstance(databases_raw, list) and databases_raw:
        if isinstance(databases_raw[0], str):
            databases = databases_raw
        elif isinstance(databases_raw[0], dict):
            databases = [d.get("name") or d.get("database") or list(d.values())[0] for d in databases_raw]
        else:
            databases = [str(d) for d in databases_raw]
    elif isinstance(databases_raw, dict):
        databases = databases_raw.get("data") or databases_raw.get("databases") or list(databases_raw.values())[0] if databases_raw else []
    else:
        databases = []

    SKIP_DBS = {"system", "information_schema", "INFORMATION_SCHEMA"}
    databases = [d for d in databases if d not in SKIP_DBS]
    print(f"   Databases found: {databases}")

    schema_lines, schema_dict = [], {}

    for db in databases:
        tables_raw = _proxy_get(f"/tables/{db}")
        if not tables_raw:
            continue
        if isinstance(tables_raw, list) and tables_raw:
            tables = tables_raw if isinstance(tables_raw[0], str) else [t.get("name") or t.get("table") or list(t.values())[0] for t in tables_raw]
        elif isinstance(tables_raw, dict):
            tables = tables_raw.get("data") or tables_raw.get("tables") or list(tables_raw.values())[0] if tables_raw else []
        else:
            tables = []

        print(f"   {db}: tables = {tables}")
        for tbl in tables:
            schema_raw = _proxy_get(f"/schema/{db}/{tbl}")
            if not schema_raw:
                schema_lines.append(f"\nTABLE: {db}.{tbl}\n  (schema unavailable)")
                continue
            if isinstance(schema_raw, list):
                cols = schema_raw
            elif isinstance(schema_raw, dict):
                cols = schema_raw.get("columns") or schema_raw.get("data") or schema_raw.get("schema") or [schema_raw]
            else:
                cols = []

            schema_dict[f"{db}.{tbl}"] = cols
            col_lines = []
            for col in cols:
                if isinstance(col, dict):
                    col_name = col.get("name") or col.get("column_name") or col.get("Field") or list(col.keys())[0]
                    col_type = col.get("type") or col.get("data_type") or col.get("Type") or ""
                    col_comment = col.get("comment") or col.get("Comment") or ""
                    col_lines.append(f"  {col_name:<35} {col_type}" + (f"  — {col_comment}" if col_comment else ""))
                else:
                    col_lines.append(f"  {col}")
            schema_lines.append(f"\nTABLE: {db}.{tbl}")
            schema_lines.extend(col_lines)

    _LIVE_SCHEMA  = schema_dict
    _SCHEMA_BLOCK = "\n".join(schema_lines) if schema_lines else "No tables found."
    print(f"✅ Schema loaded: {list(schema_dict.keys())}")
    return _SCHEMA_BLOCK


# =============================================================================
# System prompt
# =============================================================================
def _build_system_prompt() -> str:
    compact_lines = []
    for table_key, cols in _LIVE_SCHEMA.items():
        col_parts = []
        for col in cols:
            if isinstance(col, dict):
                name = col.get("name") or col.get("column_name") or col.get("Field") or list(col.keys())[0]
                typ  = col.get("type") or col.get("data_type") or col.get("Type") or ""
                col_parts.append(f"{name}:{typ}")
            else:
                col_parts.append(str(col))
        compact_lines.append(f"{table_key}({', '.join(col_parts)})")

    schema = "\n".join(compact_lines) or "Schema not yet loaded."
    if len(schema) > 20000:
        schema = schema[:20000] + "\n[schema truncated]"

    return f"""
You are DIUD (Decision Intelligence Using Data) — a conversational data assistant.
You have LIVE access to a ClickHouse database via the query_clickhouse tool.

=================================================================
GREETING RULE — HIGHEST PRIORITY
=================================================================
If the user's message is ONLY a greeting (hi, hey, hello, good morning, etc.),
respond with EXACTLY:
"Hey, I'm DIUD, your data intelligence agent to help you analyse
the live ClickHouse or Web data. How may I help you?"
No bullet points, no extras. This overrides everything.

=================================================================
EXPORT INTENT RULE
=================================================================
When the user asks to export, download, or get a list/CSV/PDF of results
from a PREVIOUS query in this conversation (e.g. "give me those 256 deals",
"export the list", "download this as CSV", "I need those deals in PDF"),
respond with this EXACT marker on a line by itself:

__EXPORT_INTENT__

Then on the next line, write a friendly confirmation message like:
"Sure! I'm exporting all [N] deals from the previous query to your chosen format."

Do NOT re-run the query. Do NOT ask which format. The export panel
will handle format selection and will re-use the already-stored query result.

=================================================================
CLICKHOUSE DIRECT ACCESS
=================================================================
You have a tool called query_clickhouse.
Use it for any question about pipeline deals, AEs, regions, industries,
stages, win/loss, competitors, conversions, or any metric not already
in the conversation context.

If the tool returns DATABASE CONNECTION FAILED, relay it to the user.

=================================================================
DUPLICATE RECORD EXCLUSION — ALWAYS APPLY
=================================================================
1. hs_analytics tables: ALWAYS use FINAL keyword
2. Aggregations: always countDistinct(), never count()
3. Association tables: DISTINCT in subquery
4. Targets table: always GROUP BY + SUM

=================================================================
TABLES
=================================================================
── TABLE 1: hs_analytics.deals ─────────────────────────────────
ALWAYS use FINAL. Key columns: deal_id, deal_name, deal_owner, deal_stage,
deal_type, pipeline, amount, region, deal_source_rollup, kore_primary_industry,
account_priority_level, create_date, close_date, became_5_deal_date,
became_10_deal_date, became_20_deal_date, became_30_deal_date,
became_40_deal_date, became_60_deal_date, became_75_deal_date

── TABLE 2: hs_analytics.owners (FINAL) ─────────────────────────
id, firstName, lastName, email

── TABLE 3: hs_analytics.companies (FINAL) ──────────────────────
company_id, name, domain, industry, country, city

── TABLE 4: hs_analytics.contacts (FINAL) ───────────────────────
contact_id, email, first_name, last_name, company_name, company_priority,
region, original_source, lead_status, lifecycle_stage,
date_entered_marketing_qualified_lead_lifecycle_stage_pipeline

── TABLE 5: kore_ai_hubspot.gs_DealContactAssociation ───────────
contact_id, deal_id

── TABLE 6: kore_ai_hubspot.gs_marketing_targets ────────────────
fy, quarter, month, region, original_source, mql_target

── TABLE 7: kore_ai_hubspot.gs_deal_ids_hs ──────────────────────
deal_id_hs

=================================================================
MANDATORY BASE FILTERS (every deals query)
=================================================================
WHERE pipeline = 'default'
AND CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
    NOT IN ('Partner-Led SMB')
AND toInt64(deal_id) IN (
    SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
)

=================================================================
FISCAL YEAR
=================================================================
FY27 = Apr 2026 – Mar 2027. Default to FY27 unless specified.
Active pipeline close_date: >= '2026-04-01' AND <= '2027-03-31'
Stages: '20% - Solution','30% - Proof','40% - Proposal',
        '60% - Price Negotiation','75% - Contract Review'

=================================================================
QUERY RULES
=================================================================
1. SELECT / WITH only — no destructive SQL
2. FINAL on all hs_analytics tables
3. Apply all 3 mandatory base filters on deals
4. For LIST queries: NO LIMIT unless user says "top N" or "first N"
   Return ALL matching rows — the system handles display safely
5. countDistinct(deal_id) for unique counts
6. round(sum(amount)/1e6, 1) for $M amounts
7. Dates: toDate(LEFT(coalesce(col,'1900-01-01'),10))
8. Always tell the user the TOTAL count (e.g. "Found 256 deals")

=================================================================
REGION / SOURCE / INDUSTRY MAPPINGS (SELECT only, not WHERE)
=================================================================
Region:  japac→JAPAC, Africa→Middle East, india___sea→ISEA
Source:  Executive Outreach+Investor→Executive Outreach, BDR Outbound→BDR, Partner→Partner - Non Hyperscaler
Industry: Financial Services+Banking+Insurance→Financial Services,
          Manufacturing Discreet+Manufacturing Process+CPG→Manufacturing

CORE RULES:
- NEVER fabricate numbers. Query the DB for every metric.
- NEVER run destructive SQL.
- Answer in clean markdown with tables for data, bold for key numbers.

=================================================================
TARGET TABLES — EXACT SCHEMA, TIER LOGIC & SQL RULES
=================================================================

TARGET TIER DEFAULT RULE — CRITICAL
─────────────────────────────────────────────────────────────────
There are three target tiers: L2 (base), L1 (stretch), Committed.
DEFAULT: Always use L2 targets unless the user explicitly says L1 or Committed.

COLUMN NAMING CONVENTION (applies across ALL target tables):
  • L2 targets  → no prefix:            Amount_Target_20, Deals_Target_20
  • L1 targets  → L1_ prefix:           L1_Amount_Target_20, L1_Deals_Target_20
  • Committed   → Committed_ prefix:    Committed_Amount_Target_20, Committed_Deals_Target_20

When user says "target" or "quota" with no qualifier → use the NO-PREFIX columns (L2).
When user says "L1" or "stretch" → use L1_ columns.
When user says "committed" → use Committed_ columns.

─────────────────────────────────────────────────────────────────
TABLE T1: kore_ai_hubspot.gs_pipeline_quotas_v1
PURPOSE : Org-wide pipeline targets by region, source, and funnel stage.
USE FOR : Pipeline attainment, EOP tracking, coverage ratio, gap-to-target.
─────────────────────────────────────────────────────────────────
COLUMNS:
  FY                          String   — e.g. 'FY27'
  Quarter                     String   — e.g. 'Q1'
  Month                       String   — e.g. '2026-04'
  Monthly_Share               Float64  — month's share of quarterly target
  Quarterly_Share             Float64  — quarter's share of annual target
  Region                      String   — 'North America','EMEA','ISEA','Global'
  Regional_Share              Float64  — region's share of global target
  Source                      String   — deal source rollup
  Source_Share                Float64  — source's share of regional target

  ── L2 targets (DEFAULT — use these unless told otherwise) ──
  Amount_Target_20            Float64  — pipeline $ target at 20%+ stage
  Deals_Target_20             Float64  — deal count target at 20%+ stage
  Amount_Target_10            Float64  — pipeline $ target at 10%+ stage
  Deals_Target_10             Float64  — deal count target at 10%+ stage
  Amount_Target_5             Float64  — pipeline $ target at 5%+ stage
  Deals_Target_5              Float64  — deal count target at 5%+ stage

  ── L1 targets (use only if user says "L1" or "stretch") ──
  Amount_Target_20_L1         Float64
  Deals_Target_20_L1          Float64
  Amount_Target_10_L1         Float64
  Deals_Target_10_L1          Float64
  Amount_Target_5_L1          Float64
  Deals_Target_5_L1           Float64

  ── Committed targets (use only if user says "committed") ──
  Amount_Target_20_committed  Float64
  Deals_Target_20_committed   Float64
  Amount_Target_10_committed  Float64
  Deals_Target_10_committed   Float64
  Amount_Target_5_committed   Float64
  Deals_Target_5_committed    Float64

─────────────────────────────────────────────────────────────────
TABLE T2: kore_ai_hubspot.gs_partner_targets_region_wise
PURPOSE : Region-level partner pipeline targets by partner type.
USE FOR : Partner pipeline attainment by region, hyperscaler vs non-hyperscaler splits.
─────────────────────────────────────────────────────────────────
COLUMNS:
  FY                          String
  Quarter                     String
  Month                       String
  Region                      String
  regional_split              Float64
  Partner_Team                String   — partner team name
  Partner_Team_Type           String   — e.g. 'Hyperscaler', 'GSI/SI', 'Reseller/BPO/TSD'
  Hyperscaler_Type            String   — 'AWS', 'MSFT', or null
  Amount_PK                   Float64  — primary key amount (total partner target)

  ── L2 targets (DEFAULT) ──
  L2_Amount_Target_20         Float64
  L2_Deals_Target_20          Float64
  L2_Amount_Target_10         Float64
  L2_Deals_Target_10          Float64
  L2_Amount_Target_5          Float64
  L2_Deals_Target_5           Float64

  ── L1 targets ──
  L1_Amount_Target_20         Float64
  L1_Deals_Target_20          Float64
  L1_Amount_Target_10         Float64
  L1_Deals_Target_10          Float64
  L1_Amount_Target_5          Float64
  L1_Deals_Target_5           Float64

  ── Committed targets ──
  Committed_Amount_Target_20  Float64
  Committed_Deals_Target_20   Float64
  Committed_Amount_Target_10  Float64
  Committed_Deals_Target_10   Float64
  Committed_Amount_Target_5   Float64
  Committed_Deals_Target_5    Float64

  ── Hyperscaler C1 targets (specific to AWS / MSFT deals) ──
  MSFT_C1_Targets_20          Float64  — MSFT C1 deal count target at 20%
  AWS_C1_Targets_20           Float64  — AWS C1 deal count target at 20%
  MSFT_C1_Amount_Target_20    Float64  — MSFT C1 $ target at 20%
  AWS_C1_Amount_Target_20     Float64  — AWS C1 $ target at 20%
  MSFT_C1_Targets_10          Float64
  AWS_C1_Targets_10           Float64
  MSFT_C1_Targets_5           Float64
  AWS_C1_Targets_5            Float64

NOTE: This table has extra trailing columns from the source sheet (FK, %share, etc.)
      — ignore those; they are sheet artefacts, not queryable metrics.

─────────────────────────────────────────────────────────────────
TABLE T3: kore_ai_hubspot.gs_partner_targets_psd
PURPOSE : PSD (Partner Sales Director) level partner pipeline targets.
USE FOR : PSD quota attainment, individual PSD performance vs target.
─────────────────────────────────────────────────────────────────
COLUMNS:
  FY                          String
  Quarter                     String
  Month                       String
  Region                      String
  Partner_Team                String
  PSD                         String   — Partner Sales Director name
  Hyperscaler_Type            String   — 'AWS', 'MSFT', or null

  ── L2 targets (DEFAULT) ──
  L2_Amount_Target_20         Float64
  L2_Deals_Target_20          Float64
  L2_Amount_Target_10         Float64
  L2_Deals_Target_10          Float64
  L2_Amount_Target_5          Float64
  L2_Deals_Target_5           Float64

  ── L1 targets ──
  L1_Amount_Target_20         Float64
  L1_Deals_Target_20          Float64
  L1_Amount_Target_10         Float64
  L1_Deals_Target_10          Float64
  L1_Amount_Target_5          Float64
  L1_Deals_Target_5           Float64

  Amount_Primary_Key          Float64  — total PSD target (primary key)

  ── Committed targets ──
  Committed_Amount_Target_20  Float64
  Committed_Deals_Target_20   Float64
  Committed_Amount_Target_10  Float64
  Committed_Deals_Target_10   Float64
  Committed_Amount_Target_5   Float64
  Committed_Deals_Target_5    Float64

─────────────────────────────────────────────────────────────────
TABLE T4: kore_ai_hubspot.gs_marketing_targets
PURPOSE : Marketing MQL and pipeline targets by source.
USE FOR : MQL attainment, marketing-sourced pipeline vs target, funnel analysis.
─────────────────────────────────────────────────────────────────
COLUMNS:
  FY                          String
  Quarter                     String
  Month                       String
  Monthly_Share               Float64
  Quarterly_Share             Float64
  Region                      String
  Regional_Share              Float64
  Original_Source             String   — maps to contacts.original_source / deals.deal_source_rollup
  Source_Share                Float64

  ── L2 targets (DEFAULT) ──
  Amount_Target_20            Float64  — marketing pipeline $ at 20%+
  Deals_Target_20             Float64
  Amount_Target_10            Float64
  Deals_Target_10             Float64
  Amount_Target_5             Float64
  Deals_Target_5              Float64
  mql_target                  Float64  — MQL count target (L2 default)

  ── L1 targets ──
  L1_mql_target               Float64
  L1_Deals_Target_20          Float64
  L1_Deals_Target_10          Float64
  L1_Deals_Target_5           Float64

NOTE: No Committed tier in this table. For MQL queries, JOIN to
      hs_analytics.contacts on region + original_source + month.
      Always: GROUP BY region, Original_Source + SUM(mql_target)

─────────────────────────────────────────────────────────────────
TABLE T5: kore_ai_hubspot.gs_closed_won_quotas
PURPOSE : Closed Won revenue quotas by AE, manager, region.
USE FOR : CW attainment %, forecast vs actual, AE-level quota tracking.
─────────────────────────────────────────────────────────────────
COLUMNS:
  fy                          String   — e.g. 'FY27'
  quarter                     String   — e.g. 'Q1'
  month                       String   — e.g. '2026-04'
  region                      String
  ae                          String   — AE name (join to deals.deal_owner)
  role                        String   — AE role/tier
  manager                     String   — AE's manager name
  assigned_amount_quota       Float64  — quarterly CW $ quota for this AE
  assigned_deals_quota        Float64  — quarterly CW deal count quota
  annualized_amount_quota     Float64  — annualized CW $ quota
  annualized_deals_quota      Float64  — annualized deal count quota

NOTE: Only one quota tier (no L1/L2/Committed split).
      To compute attainment: JOIN to deals where deal_stage = 'Closed Won'
      on ae = deal_owner, matching fy + quarter. Aggregate each side first —
      never join raw deal rows to raw quota rows before summing.

=================================================================
TARGETS SQL RULES (apply to ALL target tables)
=================================================================
1. DEFAULT TIER = L2. Only switch if user explicitly says L1 or Committed.
2. NEVER join raw deal rows directly to a target table then SUM —
   fan-out inflates targets by N× (one quota row matched N deal rows).
3. CORRECT PATTERN — two independent CTEs, then combine:

   WITH actual AS (
     SELECT region, SUM(amount) AS achieved
     FROM hs_analytics.deals FINAL
     WHERE 
     GROUP BY region
   ),
   target AS (
     SELECT region, SUM(Amount_Target_20) AS target_val   -- L2 default
     FROM kore_ai_hubspot.gs_pipeline_quotas_v1
     WHERE 
     GROUP BY region
   )
   SELECT actual.achieved, target.target_val,
          round(actual.achieved / target.target_val * 100, 1) AS attainment_pct
   FROM actual JOIN target USING (region)

4. Always GROUP BY + SUM on target tables — never read a single raw row.
5. Match period grain: quarterly actuals → quarterly target filter.
6. ATTAINMENT = round(actual / target * 100, 1)
7. COVERAGE   = round(pipeline / revenue_target, 1)
8. For partner tables: filter Partner_Team_Type to isolate
   Hyperscaler vs GSI/SI vs Reseller/BPO/TSD as needed.

=================================================================
DASHBOARD DEFINITIONS — CONTEXT & KPI LOGIC
=================================================================
When a user asks about a specific dashboard (EOP, Exec KPI, CS,
or Global Pipeline Governance), apply the correct logic below.

── DASHBOARD 1: EOP (End-of-Period) DASHBOARD ──────────────────
PURPOSE: Tracks pipeline health and attainment against EOP targets
at the end of each fiscal quarter.

KEY METRICS:
  • EOP Pipeline Value — total amount of active deals within the
    EOP date window. Source: hs_analytics.deals FINAL
  • EOP Target — from kore_ai_hubspot.gs_pipeline_quotas_v1
  • EOP Attainment % — EOP Pipeline ÷ EOP Target × 100
  • Stage-wise EOP breakdown — pipeline bucketed by deal_stage
  • Region-wise EOP — pipeline grouped by region

FILTERS TO APPLY:
  • Mandatory base filters on deals
  • close_date within current quarter end window
  • deal_stage IN active stages (20%–75%)
  • pipeline = 'default'

TYPICAL QUERIES:
  "What is our EOP pipeline vs target for Q2 FY27?"
  "Show EOP attainment by region"
  "Gap to EOP target this quarter"

── DASHBOARD 2: EXEC KPI DASHBOARD ─────────────────────────────
PURPOSE: Senior leadership view of pipeline performance, win rates,
and revenue attainment across all regions.

KEY METRICS:
  • Total Active Pipeline ($M) — sum(amount) on active deals
  • Closed Won ($M) — sum(amount) where deal_stage = 'Closed Won'
  • Closed Won Attainment % — Closed Won ÷ gs_closed_won_quotas × 100
  • Win Rate % — Closed Won deals ÷ (Closed Won + Closed Lost) × 100
  • Pipeline Coverage — Active Pipeline ÷ Revenue Target
  • New Logo Count — countDistinct(deal_id) where deal_type = 'New Logo'
  • ACV Weighted Pipeline — (stage_probability × amount) summed

FILTERS TO APPLY:
  • All mandatory base filters
  • FY27 date range on close_date
  • Exclude deal_stage IN ('Closed Won','Closed Lost') for active pipeline

TYPICAL QUERIES:
  "Executive KPI summary for FY27"
  "Closed Won attainment vs quota by region"
  "Win rate trend by quarter"

── DASHBOARD 3: CS (Customer Success) DASHBOARD ────────────────
PURPOSE: Tracks existing customer pipeline — renewals, upsells,
expansions — and CS team performance.

KEY METRICS:
  • Renewal Pipeline ($M) — deals where deal_type LIKE '%Renewal%'
  • Upsell / Expansion Pipeline ($M) — deal_type LIKE '%Upsell%'
    or deal_type LIKE '%Expansion%'
  • Renewal Win Rate % — Closed Won renewals ÷ total renewals × 100
  • Net Revenue Retention (NRR) — (Renewals + Upsells) ÷ Base ARR
  • At-Risk Deals — active deals with stale last_contacted date
  • CS AE Performance — pipeline / closed won by owner filtered to
    CS team (join to kore_ai_hubspot.gs_Teams on hubspot_team)

FILTERS TO APPLY:
  • All mandatory base filters
  • deal_type IN ('Renewal','Upsell','Expansion') or similar values
  • FY27 date range

TYPICAL QUERIES:
  "CS renewal pipeline for FY27"
  "Upsell attainment by AE"
  "At-risk renewals this quarter"

── DASHBOARD 4: GLOBAL PIPELINE GOVERNANCE DASHBOARD ───────────
PURPOSE: Executive governance view comparing pipeline across all
regions, sources, and partner types against global targets.

KEY METRICS:
  • Global Pipeline by Region ($M) — broken down by region + stage
  • Partner Pipeline ($M) — deals from partner sources
    (deal_source_rollup LIKE '%Partner%')
  • Partner Attainment % — vs kore_ai_hubspot.gs_partner_targets_region_wise
  • Partner PSD Attainment % — vs kore_ai_hubspot.gs_partner_targets_psd
  • Pipeline Coverage Ratio — by region vs gs_pipeline_quotas_v1
  • Closed Won Governance — actual vs gs_closed_won_quotas by region/quarter
  • Marketing Sourced Pipeline — deals from marketing sources
    vs gs_marketing_targets

FILTERS TO APPLY:
  • All mandatory base filters
  • FY27 date range
  • Appropriate partner source filters for partner metrics

TYPICAL QUERIES:
  "Global pipeline governance report for FY27"
  "Partner pipeline attainment by region"
  "Closed Won vs quota by quarter"
  "Marketing sourced pipeline vs targets"

DASHBOARD SELECTION RULE:
If the user mentions a specific dashboard by name, apply its metric
definitions and target table references automatically. If unclear,
ask the user which dashboard context they want.
"""

_SYSTEM_PROMPT = _build_system_prompt()

# =============================================================================
# ClickHouse query runner — stores full result in session, returns display text
# =============================================================================
def run_clickhouse_query(sql: str, session_id: Optional[str] = None) -> str:
    """
    Execute SQL against ClickHouse.
    - Stores the FULL result set in _SESSION_STORE[session_id] (no row cap).
    - Returns a chat-display string capped at CHAT_DISPLAY_LIMIT rows.
    The export layer reads from the session store and gets all rows.
    """
    base_url = _base_url()
    token    = _token()

    if not base_url:
        return "DATABASE CONNECTION FAILED: CLICKHOUSE_API_URL is not set."
    if not token:
        return "DATABASE CONNECTION FAILED: CLICKHOUSE_API_TOKEN is not set."

    stripped = sql.strip().upper()
    if not (stripped.startswith("SELECT") or stripped.startswith("WITH")):
        return "ERROR: Only SELECT/WITH queries are permitted."
    for kw in FORBIDDEN_KEYWORDS:
        if re.search(rf'\b{kw}\b', stripped):
            return f"ERROR: Forbidden keyword: {kw}"

    print(f"🔍 SQL (session={session_id}) → {sql[:200]}")

    try:
        resp = httpx.post(
            f"{base_url}/query",
            headers=_auth_headers(),
            json={"query": sql},
            timeout=60,
        )

        if resp.status_code == 401:
            return "DATABASE CONNECTION FAILED: 401 Unauthorized."
        if resp.status_code == 403:
            return "DATABASE CONNECTION FAILED: 403 Forbidden."
        if resp.status_code == 422:
            return f"ERROR: Proxy rejected query (422): {resp.text[:400]}"
        if resp.status_code == 500:
            return f"DATABASE ERROR: HTTP 500: {resp.text[:400]}"
        if resp.status_code != 200:
            return f"DATABASE ERROR: HTTP {resp.status_code} — {resp.text[:300]}"

        payload = resp.json()

        # Normalise to list of rows
        if isinstance(payload, list):
            rows = payload
        elif isinstance(payload, dict):
            rows = payload.get("data") or payload.get("rows") or payload.get("result") or payload.get("results")
            api_columns = payload.get("columns") or payload.get("meta") or payload.get("column_names")
            if rows is None:
                return json.dumps(payload, indent=2, default=str)[:3000]
        else:
            return f"Unexpected response type: {type(payload)}"

        if not rows:
            return "Query returned 0 rows."

        # Normalise rows to list-of-dicts
        if isinstance(rows[0], dict):
            columns = list(rows[0].keys())
            norm_rows = rows
        else:
            # List-of-lists — use index keys
            if api_columns and len(api_columns) == len(rows[0]):
                columns = [c["name"] if isinstance(c, dict) else c for c in api_columns]
            else:
                columns = [f"col_{i}" for i in range(len(rows[0]))]
            norm_rows = [dict(zip(columns, r)) for r in rows]

        total_rows = len(norm_rows)

        # ── Store FULL result in session (no row cap) ──────────────────
        if session_id:
            _store_result(session_id, QueryResult(
                sql          = sql,
                columns      = columns,
                rows         = norm_rows,   # ALL rows
                total_rows   = total_rows,
                captured_at  = datetime.utcnow().isoformat() + "Z",
                filters_applied = _extract_filters_from_sql(sql),
            ))

        # ── Build chat display (capped at 100 rows for readability) ───
        CHAT_DISPLAY_LIMIT = 100
        header = " | ".join(columns)
        lines  = [header, "-" * min(len(header), 140)]
        for row in norm_rows[:CHAT_DISPLAY_LIMIT]:
            lines.append(" | ".join(str(row.get(c, "")) for c in columns))

        if total_rows > CHAT_DISPLAY_LIMIT:
            lines.append(
                f"\n📊 **Showing {CHAT_DISPLAY_LIMIT} of {total_rows} rows.** "
                f"Say **\"export these deals as CSV\"** or **\"export as PDF\"** "
                f"to download all {total_rows} records."
            )

        result = "\n".join(lines)
        print(f"   ✅ {total_rows} rows returned. Session store updated.")
        return result

    except httpx.ConnectError as e:
        return f"DATABASE CONNECTION FAILED: Could not reach {base_url}. {e}"
    except httpx.TimeoutException:
        return "DATABASE CONNECTION FAILED: Query timed out after 60 seconds."
    except Exception as exc:
        traceback.print_exc()
        return f"DATABASE CONNECTION FAILED: {type(exc).__name__}: {exc}"


def _extract_filters_from_sql(sql: str) -> str:
    """Extract a human-readable summary of WHERE filters from SQL."""
    sql_upper = sql.upper()
    filters = []
    if "PIPELINE = 'DEFAULT'" in sql_upper:
        filters.append("Pipeline: default")
    if "BECAME_20_DEAL_DATE" in sql_upper:
        filters.append("Cohort: 20% qualified deals")
    if "BECAME_5_DEAL_DATE" in sql_upper:
        filters.append("Cohort: 5% IQM deals")
    if "CLOSE_DATE" in sql_upper and "2026-04-01" in sql:
        filters.append("FY27 active pipeline")
    if "DEAL_STAGE" in sql_upper:
        m = re.search(r"deal_stage\s+IN\s*\(([^)]+)\)", sql, re.IGNORECASE)
        if m:
            filters.append(f"Stage filter: {m.group(1)[:60]}")
    if "REGION" in sql_upper:
        m = re.search(r"region\s*=\s*'([^']+)'", sql, re.IGNORECASE)
        if m:
            filters.append(f"Region: {m.group(1)}")
    return "; ".join(filters) if filters else "Standard base filters applied"


# =============================================================================
# Startup
# =============================================================================
@app.on_event("startup")
async def on_startup():
    global _SYSTEM_PROMPT
    discover_schema()
    _SYSTEM_PROMPT = _build_system_prompt()
    print("🚀 DIUD v4 started — session-store export enabled.")


# =============================================================================
# Claude tool definition
# =============================================================================
_QUERY_TOOL = {
    "name": "query_clickhouse",
    "description": (
        "Execute a SELECT query against ClickHouse. "
        "Use for deal pipeline, AE performance, win/loss, regions, stages, MQL metrics. "
        "ALWAYS use fully-qualified table names. "
        "Relay DATABASE CONNECTION FAILED errors directly to the user."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "sql": {
                "type": "string",
                "description": (
                    "Valid ClickHouse SELECT or WITH query. "
                    "For deal LIST queries: NO LIMIT unless user asks for 'top N'. "
                    "Return all matching rows — the system displays them safely."
                ),
            }
        },
        "required": ["sql"],
    },
}

# =============================================================================
# Pydantic models
# =============================================================================
class ChatMessage(BaseModel):
    role: Literal["user", "assistant"]
    content: str

class ChatRequest(BaseModel):
    message: str
    history: List[ChatMessage] = []
    session_id: Optional[str] = None

class ExportPreviewRequest(BaseModel):
    conversation: List[ChatMessage] = []
    title: str = "Pipeline Intelligence Report"
    export_type: Literal["pdf", "pptx"] = "pdf"
    detail_level: Literal["summary", "detailed"] = "detailed"
    session_id: Optional[str] = None

class ExportDownloadRequest(BaseModel):
    format: Literal["pdf", "pptx", "csv"]
    content: Optional[str] = None    # markdown content (pdf/pptx)
    title: str = "Pipeline Intelligence Report"
    session_id: Optional[str] = None


# =============================================================================
# Claude tool loop — passes session_id so query runner can store results
# =============================================================================
def _extract_text(content_blocks) -> str:
    return "\n".join(
        b.text for b in content_blocks if hasattr(b, "text") and b.text
    ).strip()


def _call_claude(messages: list, max_tokens: int = 2048, session_id: Optional[str] = None) -> str:
    """Run Claude with query_clickhouse tool. Up to 5 tool rounds."""

    # Strip tool_use/tool_result blocks from history (can't replay them)
    safe_messages = []
    for m in messages:
        content = m.get("content", "")
        if isinstance(content, list):
            text_parts = [
                b.get("text", "") if isinstance(b, dict) else (b.text if hasattr(b, "text") else "")
                for b in content if (isinstance(b, dict) and b.get("type") == "text")
                   or (hasattr(b, "type") and b.type == "text")
            ]
            text = "\n".join(t for t in text_parts if t).strip()
            if text:
                safe_messages.append({"role": m["role"], "content": text})
        else:
            safe_messages.append({"role": m["role"], "content": content})

    response = _ai_client.messages.create(
        model=_CLAUDE_MODEL,
        system=_SYSTEM_PROMPT,
        messages=safe_messages,
        tools=[_QUERY_TOOL],
        temperature=0,
        max_tokens=max_tokens,
    )

    MAX_ROUNDS = 8
    last_error = None

    for round_num in range(MAX_ROUNDS):
        if response.stop_reason != "tool_use":
            break

        tool_block = next((b for b in response.content if b.type == "tool_use"), None)
        if not tool_block:
            break

        sql          = tool_block.input.get("sql", "")
        query_result = run_clickhouse_query(sql, session_id=session_id)
        is_error     = any(query_result.startswith(p) for p in [
            "DATABASE CONNECTION FAILED", "ERROR:", "DATABASE ERROR:"
        ])
        if is_error:
            last_error = query_result

        safe_messages = safe_messages + [
            {"role": "assistant", "content": response.content},
            {"role": "user", "content": [{
                "type":        "tool_result",
                "tool_use_id": tool_block.id,
                "content":     query_result,
                "is_error":    is_error,
            }]},
        ]

        is_last_round = (round_num == MAX_ROUNDS - 1)
        response = _ai_client.messages.create(
            model=_CLAUDE_MODEL,
            system=_SYSTEM_PROMPT,
            messages=safe_messages,
            # On the final round, withhold tools so Claude is forced to
            # summarize in text instead of issuing yet another tool call.
            tools=[] if is_last_round else [_QUERY_TOOL],
            temperature=0,
            max_tokens=max_tokens,
        )

    reply = _extract_text(response.content)

    if not reply:
        # Still empty even after forcing a text-only round. This means the
        # model returned no text at all (rare) — surface the real cause
        # instead of a generic "connectivity" message.
        if last_error:
            reply = (
                "⚠️ I couldn't complete this query. The last database error was:\n\n"
                f"`{last_error[:400]}`\n\n"
                "Could you rephrase the question, or check **/debug/db** if this persists?"
            )
        else:
            reply = (
                "⚠️ I wasn't able to finish answering this in time — it may need "
                "a more specific or simpler question. Could you try rephrasing it "
                "(e.g. break it into smaller asks)?"
            )
    return reply


# =============================================================================
# Routes — chat
# =============================================================================
@app.get("/", response_class=HTMLResponse)
def root():
    with open("chat.html", "r") as f:
        return HTMLResponse(content=f.read())

@app.get("/logo.png")
def serve_logo():
    return FileResponse("logo.png", media_type="image/png")


@app.get("/debug/db")
def debug_db():
    base_url = _base_url()
    token    = _token()
    config = {
        "CLICKHOUSE_API_URL":   base_url or "❌ NOT SET",
        "CLICKHOUSE_API_TOKEN": f"✅ set ({len(token)} chars)" if token else "❌ NOT SET",
    }
    if not base_url or not token:
        return {"status": "MISCONFIGURED", "config": config}

    tests = {}
    try:
        r = httpx.get(base_url, timeout=10)
        tests["GET /"] = {"status": r.status_code}
    except Exception as e:
        tests["GET /"] = {"error": str(e)}

    ping = run_clickhouse_query("SELECT 1 AS ping")
    query_ok = not any(ping.startswith(p) for p in ["DATABASE CONNECTION FAILED", "ERROR:", "DATABASE ERROR:"])
    tests["SELECT 1"] = {"ok": query_ok, "result": ping[:100]}

    return {
        "status":            "OK" if query_ok else "FAILED",
        "config":            config,
        "discovered_tables": list(_LIVE_SCHEMA.keys()),
        "active_sessions":   len(_SESSION_STORE),
        "tests":             tests,
    }


@app.post("/refresh-schema")
def refresh_schema():
    global _SYSTEM_PROMPT
    schema = discover_schema()
    _SYSTEM_PROMPT = _build_system_prompt()
    return {"status": "refreshed", "tables": list(_LIVE_SCHEMA.keys())}


@app.post("/chat")
def chat(payload: ChatRequest):
    messages = [{"role": m.role, "content": m.content} for m in payload.history]
    messages.append({"role": "user", "content": payload.message})
    print(f"💬 [chat] session={payload.session_id} msg={payload.message[:80]}")

    try:
        reply = _call_claude(messages, session_id=payload.session_id)
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Claude error: {exc}")

    # Surface whether there is a stored dataset for this session
    has_dataset = payload.session_id is not None and payload.session_id in _SESSION_STORE
    stored = _SESSION_STORE.get(payload.session_id) if payload.session_id else None

    return {
        "reply":        reply,
        "has_dataset":  has_dataset,
        "dataset_rows": stored.total_rows if stored else 0,
        "export_intent": "__EXPORT_INTENT__" in reply,
    }


# =============================================================================
# Retry — re-runs the last user message with fresh LLM call
# =============================================================================

class RetryRequest(BaseModel):
    """
    Re-run the most recent user message with a fresh LLM call.

    history    : full conversation UP TO AND INCLUDING the last user message.
                 The last item must be role='user'. Any prior assistant reply
                 for that turn is intentionally excluded so the model generates
                 a new response.
    session_id : existing session — query results from previous turns are
                 preserved in the session store so exports still work.
    """
    history:    List[ChatMessage] = []
    session_id: Optional[str] = None


@app.post("/chat/retry")
def chat_retry(payload: RetryRequest):
    """
    Regenerate the last assistant response without adding a new user message.

    Steps:
    1. Validate that the last history entry is a user message.
    2. Remove the last assistant message if present (prevents duplicate in history).
    3. Call Claude fresh with the same conversation context.
    4. Return the new reply with the same response shape as /chat.

    This preserves:
    - Full conversation context (all prior turns)
    - Session store (dataset / query results)
    - Dashboard context embedded in history
    """
    if not payload.history:
        raise HTTPException(status_code=400, detail="history must not be empty for retry.")

    # Walk backwards: find the last user message and strip any trailing assistant
    # message so we don't send the old answer as context for the regeneration.
    clean_history = list(payload.history)

    # Remove trailing assistant message(s) — they are the stale response being retried
    while clean_history and clean_history[-1].role == "assistant":
        clean_history.pop()

    if not clean_history or clean_history[-1].role != "user":
        raise HTTPException(
            status_code=400,
            detail="No user message found to retry. history must end with a user turn."
        )

    last_user_msg = clean_history[-1].content
    prior_history = clean_history[:-1]   # everything before the last user message

    print(f"🔄 [retry] session={payload.session_id} retrying: {last_user_msg[:80]}")

    messages = [{"role": m.role, "content": m.content} for m in prior_history]
    messages.append({"role": "user", "content": last_user_msg})

    try:
        reply = _call_claude(messages, session_id=payload.session_id)
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Claude error on retry: {exc}")

    has_dataset = payload.session_id is not None and payload.session_id in _SESSION_STORE
    stored      = _SESSION_STORE.get(payload.session_id) if payload.session_id else None

    return {
        "reply":         reply,
        "has_dataset":   has_dataset,
        "dataset_rows":  stored.total_rows if stored else 0,
        "export_intent": "__EXPORT_INTENT__" in reply,
        "retried":       True,
    }


# =============================================================================
# Session info — lets the frontend know what's stored
# =============================================================================
@app.get("/session/{session_id}/dataset-info")
def session_dataset_info(session_id: str):
    result = _get_result(session_id)
    if not result:
        return {"has_dataset": False}
    return {
        "has_dataset":     True,
        "total_rows":      result.total_rows,
        "columns":         result.columns,
        "captured_at":     result.captured_at,
        "filters_applied": result.filters_applied,
        "sql_preview":     result.sql[:300],
    }


# =============================================================================
# Export preview — generates AI narrative + injects full table
# =============================================================================
@app.post("/export/preview")
async def export_preview(req: ExportPreviewRequest):
    if not req.conversation:
        raise HTTPException(status_code=400, detail="No conversation to export.")

    print(f"📄 [export/preview] session={req.session_id} type={req.export_type}")

    # Fetch stored dataset (may be None for summary-only exports)
    stored = _get_result(req.session_id) if req.session_id else None

    try:
        ai_content = _generate_export_content(
            conversation   = req.conversation,
            title          = req.title,
            export_type    = req.export_type,
            detail_level   = req.detail_level,
            stored_dataset = stored,
        )
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Content generation error: {exc}")

    return {
        "content":       ai_content,
        "title":         req.title,
        "export_type":   req.export_type,
        "word_count":    len(ai_content.split()),
        "generated_at":  date.today().isoformat(),
        "total_rows":    stored.total_rows if stored else 0,
        "filters":       stored.filters_applied if stored else "",
    }


# =============================================================================
# Export download — PDF, PPTX, or CSV, all from session store
# =============================================================================
@app.post("/export/download")
async def export_download(req: ExportDownloadRequest):
    print(f"⬇️  [export/download] format={req.format} session={req.session_id}")

    # ── CSV: purely from session store, no AI involvement ────────────────
    if req.format == "csv":
        stored = _get_result(req.session_id) if req.session_id else None
        if not stored:
            raise HTTPException(
                status_code=404,
                detail=(
                    "No query result found for this session. "
                    "Ask a deal-list question first, then export."
                ),
            )
        csv_bytes = _build_csv(stored)
        return StreamingResponse(
            io.BytesIO(csv_bytes),
            media_type="text/csv",
            headers={
                "Content-Disposition": f'attachment; filename="{_safe_filename(req.title)}.csv"',
                "X-Total-Rows": str(stored.total_rows),
            },
        )

    # ── PDF / PPTX: from pre-generated content (passed from preview step) ─
    if not req.content:
        raise HTTPException(status_code=400, detail="content is required for PDF/PPTX export.")

    try:
        if req.format == "pdf":
            file_bytes = _build_pdf(req.title, req.content)
            media_type = "application/pdf"
            ext        = "pdf"
        else:
            file_bytes = _build_pptx(req.title, req.content)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ext        = "pptx"

        return StreamingResponse(
            io.BytesIO(file_bytes),
            media_type=media_type,
            headers={"Content-Disposition": f'attachment; filename="{_safe_filename(req.title)}.{ext}"'},
        )
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"File generation error: {exc}")


# =============================================================================
# Helpers
# =============================================================================
def _safe_filename(title: str) -> str:
    return re.sub(r'[^\w\-]', '_', title)[:60]


def _strip_md(t: str) -> str:
    t = re.sub(r'\*{1,3}(.*?)\*{1,3}', r'\1', t)
    t = re.sub(r'#{1,6}\s*', '', t)
    t = re.sub(r'^[\-\*•]\s*', '', t, flags=re.M)
    t = re.sub(r'`(.*?)`', r'\1', t)
    t = t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return t.strip()


# =============================================================================
# CSV builder — generates UTF-8 CSV from session store (all rows, no cap)
# =============================================================================
def _build_csv(stored: QueryResult) -> bytes:
    buf = io.StringIO()

    # Metadata header
    buf.write(f"# Title: {stored.sql[:80]}\n")
    buf.write(f"# Generated: {date.today().isoformat()}\n")
    buf.write(f"# Total Records: {stored.total_rows}\n")
    buf.write(f"# Filters: {stored.filters_applied}\n")
    buf.write(f"# Captured at: {stored.captured_at}\n")
    buf.write("#\n")

    writer = csv.DictWriter(
        buf,
        fieldnames     = stored.columns,
        extrasaction   = "ignore",
        lineterminator = "\n",
    )
    writer.writeheader()
    for row in stored.rows:   # ALL rows from session store
        writer.writerow(row)

    return buf.getvalue().encode("utf-8")


# =============================================================================
# Export content generation
# =============================================================================
def _generate_export_content(
    conversation:    List[ChatMessage],
    title:           str,
    export_type:     str,
    detail_level:    str = "detailed",
    stored_dataset:  Optional[QueryResult] = None,
) -> str:
    """
    Claude writes the narrative (summary, insights, recommendations).
    The full deal table (if any) is injected directly from the session store —
    not reconstructed from chat text, not subject to any token/row limits.
    """

    # Clean conversation: strip any old embedded JSON blocks from previous iterations
    conv_text = "\n\n".join(
        f"{'USER' if m.role == 'user' else 'DIUD AGENT'}: {m.content}"
        for m in conversation
    )

    format_hint = (
        "Format as a PowerPoint: use SLIDE: <title> for each slide, then bullet points."
        if export_type == "pptx"
        else "Format as a professional PDF report: ## section headers, narrative prose, tables."
    )
    detail_hint = (
        "Include all metrics and insights. The full deal table will be appended automatically — "
        "just write a [DEAL_TABLE_PLACEHOLDER] marker where it should appear."
        if detail_level == "detailed"
        else "Executive summary only — key metrics and top insights, no raw deal list."
    )

    dataset_hint = ""
    if stored_dataset:
        dataset_hint = (
            f"\n\nDATASET CONTEXT: The query returned {stored_dataset.total_rows} records "
            f"with columns: {', '.join(stored_dataset.columns[:12])}. "
            f"Filters: {stored_dataset.filters_applied}. "
            f"The complete table will be injected at [DEAL_TABLE_PLACEHOLDER]."
        )

    prompt = f"""You are preparing a professional {export_type.upper()} export document.

CONVERSATION:
{conv_text}
{dataset_hint}

TASK: Create "{title}"

{format_hint}
{detail_hint}

REQUIREMENTS:
- Executive summary at the start with key numbers
- Logical sections: summary, key metrics, insights, recommendations
- If this is a deal list export, include [DEAL_TABLE_PLACEHOLDER] where the full table belongs
- Bold key numbers; clean professional tone
- Today: {date.today().strftime('%B %d, %Y')}

Generate the document now:"""

    response = _ai_client.messages.create(
        model   = _CLAUDE_MODEL,
        system  = "You are a professional business report writer. Generate clean, well-structured documents.",
        messages= [{"role": "user", "content": prompt}],
        temperature = 0,
        max_tokens  = 4096,
    )
    ai_text = _extract_text(response.content)

    # ── Inject full deal table from session store ──────────────────────────
    # Claude is NOT trusted to reproduce the table — we do it ourselves.
    if stored_dataset and stored_dataset.total_rows > 0:
        table_md  = _rows_to_markdown_table(stored_dataset)
        meta_line = (
            f"**Total records:** {stored_dataset.total_rows:,} | "
            f"**Filters:** {stored_dataset.filters_applied} | "
            f"**Exported:** {date.today().strftime('%B %d, %Y')}"
        )
        full_section = f"\n\n## Deal List ({stored_dataset.total_rows:,} records)\n\n{meta_line}\n\n{table_md}"

        if "[DEAL_TABLE_PLACEHOLDER]" in ai_text:
            ai_text = ai_text.replace("[DEAL_TABLE_PLACEHOLDER]", full_section)
        else:
            ai_text = ai_text.rstrip() + full_section

    return ai_text


def _rows_to_markdown_table(stored: QueryResult) -> str:
    """Convert session-store rows to a markdown table (all rows, no cap)."""
    if not stored.rows:
        return "_No data._"

    cols   = stored.columns
    header = "| " + " | ".join(cols) + " |"
    sep    = "| " + " | ".join("---" for _ in cols) + " |"
    lines  = [header, sep]

    for row in stored.rows:   # all rows from session store
        cells = [str(row.get(c, "")).replace("|", "\\|") for c in cols]
        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


# =============================================================================
# PDF Builder
# =============================================================================
_C_NAVY  = colors.HexColor("#0D1B3E")
_C_BLUE  = colors.HexColor("#1565C0")
_C_WHITE = colors.white
_C_BG    = colors.HexColor("#F7F9FC")
_C_TXT   = colors.HexColor("#1E293B")
_C_DIM   = colors.HexColor("#94A3B8")

_SECTION_COLORS = {
    "executive": colors.HexColor("#0D1B3E"),
    "pipeline":  colors.HexColor("#1565C0"),
    "metric":    colors.HexColor("#004D40"),
    "deal":      colors.HexColor("#1565C0"),
    "regional":  colors.HexColor("#BF360C"),
    "win":       colors.HexColor("#B71C1C"),
    "loss":      colors.HexColor("#B71C1C"),
    "recommend": colors.HexColor("#1B5E20"),
    "summary":   colors.HexColor("#0D1B3E"),
    "analysis":  colors.HexColor("#1565C0"),
    "overview":  colors.HexColor("#004D40"),
    "insight":   colors.HexColor("#1B5E20"),
}

PW, PH = A4
_ML = _MR = 0.6 * inch
_MT = 0.45 * inch
_MB = 0.40 * inch
_HDR_H = 44
_FTR_H = 20
_CW    = PW - _ML - _MR


def _pdf_styles():
    return {
        "Cover_Title": ParagraphStyle("Cover_Title", fontSize=26, leading=32,
            textColor=_C_WHITE, fontName="Helvetica-Bold", spaceAfter=8),
        "Cover_Sub": ParagraphStyle("Cover_Sub", fontSize=13, leading=18,
            textColor=colors.HexColor("#B0BEC5"), fontName="Helvetica"),
        "Section_H": ParagraphStyle("Section_H", fontSize=11, leading=15,
            textColor=_C_WHITE, fontName="Helvetica-Bold"),
        "Body": ParagraphStyle("Body", fontSize=9, leading=14, textColor=_C_TXT,
            fontName="Helvetica", spaceAfter=4),
        "Bullet": ParagraphStyle("Bullet", fontSize=9, leading=14, textColor=_C_TXT,
            fontName="Helvetica", leftIndent=12, firstLineIndent=-8, spaceAfter=3),
        "H2": ParagraphStyle("H2", fontSize=11, leading=15, textColor=_C_NAVY,
            fontName="Helvetica-Bold", spaceBefore=10, spaceAfter=4),
        "H3": ParagraphStyle("H3", fontSize=9, leading=13, textColor=_C_BLUE,
            fontName="Helvetica-Bold", spaceBefore=6, spaceAfter=2),
        "TH": ParagraphStyle("TH", fontSize=7, leading=9, textColor=_C_WHITE,
            fontName="Helvetica-Bold"),
        "TD": ParagraphStyle("TD", fontSize=7, leading=9, textColor=_C_TXT,
            fontName="Helvetica"),
    }


def _parse_sections(text: str):
    parts = re.split(r'^##\s+', text, flags=re.MULTILINE)
    return [
        (lines[0].strip(), lines[1].strip() if len(lines) > 1 else "")
        for part in parts if part.strip()
        for lines in [part.strip().split("\n", 1)]
    ]


def _build_pdf(title: str, report_text: str) -> bytes:
    """
    Build a PDF from markdown content.
    Markdown tables (| col | col |) are rendered as native ReportLab tables
    so that deal lists with many columns stay readable.
    """
    buf     = io.BytesIO()
    styles  = _pdf_styles()
    sections = _parse_sections(report_text)

    def _on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(_C_NAVY)
        canvas.rect(0, PH - _HDR_H - _MT, PW, _HDR_H + _MT, fill=1, stroke=0)
        canvas.setFillColor(_C_WHITE)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawString(_ML, PH - _MT - 28, title)
        canvas.setFillColor(_C_BG)
        canvas.rect(0, 0, PW, _FTR_H + _MB, fill=1, stroke=0)
        canvas.setFillColor(_C_DIM)
        canvas.setFont("Helvetica", 7)
        canvas.drawCentredString(
            PW / 2, _MB + 5,
            f"DIUD Report  |  AI-Generated  |  CONFIDENTIAL  |  {date.today().strftime('%B %Y')}"
        )
        canvas.drawRightString(PW - _MR, _MB + 5, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    frame    = Frame(_ML, _MB + _FTR_H, _CW, PH - _HDR_H - _MT - _MB - _FTR_H, id="main")
    template = PageTemplate(id="main", frames=[frame], onPage=_on_page)
    doc = BaseDocTemplate(buf, pagesize=A4, leftMargin=_ML, rightMargin=_MR,
                          topMargin=_MT + _HDR_H, bottomMargin=_MB + _FTR_H)
    doc.addPageTemplates([template])

    story = [
        Spacer(1, 1.0 * inch),
        Paragraph(title, styles["Cover_Title"]),
        Paragraph(f"Generated {date.today().strftime('%B %d, %Y')}", styles["Cover_Sub"]),
        PageBreak(),
    ]

    if not sections:
        for line in report_text.split("\n"):
            line = line.strip()
            if line:
                story.append(Paragraph(_strip_md(line), styles["Body"]))
    else:
        for sec_title, sec_body in sections:
            color_key = next((k for k in _SECTION_COLORS if k in sec_title.lower()), None)
            bar_color = _SECTION_COLORS.get(color_key, _C_BLUE)
            story.append(Table(
                [[Paragraph(sec_title.upper(), styles["Section_H"])]],
                colWidths=[_CW],
                style=TableStyle([
                    ("BACKGROUND",    (0, 0), (-1, -1), bar_color),
                    ("TOPPADDING",    (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING",   (0, 0), (-1, -1), 10),
                ])
            ))
            story.append(Spacer(1, 6))

            # Detect markdown table blocks vs normal lines
            lines = sec_body.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i].strip()

                # Start of a markdown table
                if line.startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith("|"):
                        table_lines.append(lines[i].strip())
                        i += 1
                    story.append(_md_table_to_rl(table_lines, styles))
                    story.append(Spacer(1, 6))
                    continue

                if not line:
                    story.append(Spacer(1, 3))
                elif line.startswith("### "):
                    story.append(Paragraph(_strip_md(line), styles["H3"]))
                elif line.startswith("## "):
                    story.append(Paragraph(_strip_md(line), styles["H2"]))
                elif line.startswith(("- ", "* ", "• ")):
                    story.append(Paragraph("• " + _strip_md(line[2:]), styles["Bullet"]))
                else:
                    story.append(Paragraph(_strip_md(line), styles["Body"]))
                i += 1

            story.extend([Spacer(1, 12), PageBreak()])

    doc.build(story)
    return buf.getvalue()


def _md_table_to_rl(table_lines: list, styles: dict):
    """Convert markdown table lines to a ReportLab Table element."""
    data = []
    for idx, line in enumerate(table_lines):
        if "---" in line:    # separator row — skip
            continue
        cells = [c.strip().replace("\\|", "|") for c in line.strip("|").split("|")]
        if idx == 0:
            row = [Paragraph(_strip_md(c), styles["TH"]) for c in cells]
        else:
            row = [Paragraph(_strip_md(c), styles["TD"]) for c in cells]
        data.append(row)

    if not data:
        return Spacer(1, 1)

    num_cols = max(len(r) for r in data)
    col_w    = _CW / max(num_cols, 1)

    tbl = Table(data, colWidths=[col_w] * num_cols, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0),  _C_NAVY),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  _C_WHITE),
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, -1), 7),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
        ("GRID",          (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
        ("TOPPADDING",    (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ("LEFTPADDING",   (0, 0), (-1, -1), 4),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
    ]))
    return tbl


# =============================================================================
# PPTX Builder
# =============================================================================
_C_NAVY_P  = RGBColor(0x0D, 0x1B, 0x3E)
_C_DNAV_P  = RGBColor(0x0A, 0x11, 0x28)
_C_BLUE_P  = RGBColor(0x1E, 0x88, 0xE5)
_C_WHITE_P = RGBColor(0xFF, 0xFF, 0xFF)
_C_LTBG_P  = RGBColor(0xF5, 0xF7, 0xFA)
_C_TXT_P   = RGBColor(0x1A, 0x1A, 0x2E)
_C_DIM_P   = RGBColor(0x88, 0x99, 0xAA)

_SLIDE_ACCENT = {
    "overview":  RGBColor(0x1E, 0x88, 0xE5),
    "pipeline":  RGBColor(0x00, 0x89, 0x7B),
    "deal":      RGBColor(0x1E, 0x88, 0xE5),
    "metric":    RGBColor(0x2E, 0x7D, 0x32),
    "regional":  RGBColor(0xBF, 0x36, 0x0C),
    "win":       RGBColor(0x2E, 0x7D, 0x32),
    "loss":      RGBColor(0xC6, 0x28, 0x28),
    "recommend": RGBColor(0x1B, 0x5E, 0x20),
    "summary":   RGBColor(0x1E, 0x88, 0xE5),
    "analysis":  RGBColor(0x00, 0x89, 0x7B),
    "insight":   RGBColor(0x1B, 0x5E, 0x20),
    "executive": RGBColor(0x0D, 0x1B, 0x3E),
}


def _pptx_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _pptx_rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    return shp

def _pptx_txt(slide, text, l, t, w, h, bold=False, size=18, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color or _C_TXT_P
    return txb

def _parse_slides(text: str):
    slides, cur_title, cur_bullets = [], None, []
    for line in text.split("\n"):
        line = line.rstrip()
        if line.startswith("SLIDE:"):
            if cur_title is not None:
                slides.append((cur_title, cur_bullets))
            cur_title, cur_bullets = line[6:].strip(), []
        elif line.startswith("- ") and cur_title:
            cur_bullets.append(line[2:].strip())
    if cur_title is not None:
        slides.append((cur_title, cur_bullets))
    return slides

def _build_pptx(title: str, slide_text: str) -> bytes:
    slides_data = _parse_slides(slide_text) or [(title, [slide_text[:400]])]
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    footer_text = f"DIUD  |  AI-Generated  |  CONFIDENTIAL  |  {date.today().strftime('%B %Y')}"
    blank = prs.slide_layouts[6]

    def _footer(s):
        _pptx_rect(s, 0, 7.1, 13.33, 0.4, _C_DNAV_P)
        _pptx_txt(s, footer_text, 0.3, 7.12, 12, 0.35, size=7, color=_C_DIM_P, align=PP_ALIGN.CENTER)

    def _accent(t):
        for k, c in _SLIDE_ACCENT.items():
            if k in t: return c
        return _C_BLUE_P

    # Cover
    cover = prs.slides.add_slide(blank)
    _pptx_bg(cover, _C_NAVY_P)
    _pptx_rect(cover, 0, 3.2, 13.33, 0.06, _C_BLUE_P)
    _pptx_txt(cover, title, 0.8, 1.6, 11.5, 1.4, bold=True, size=34, color=_C_WHITE_P)
    _pptx_txt(cover, "Deals Intelligence Report", 0.8, 3.0, 8, 0.6,
              size=15, color=RGBColor(0xB0, 0xBE, 0xC5))
    _pptx_txt(cover, f"Generated: {date.today().strftime('%B %d, %Y')}", 0.8, 3.6, 6, 0.45,
              size=12, color=RGBColor(0x78, 0x90, 0x9C))

    for i, (s_title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        ac = _accent(s_title.lower())
        _pptx_bg(slide, _C_LTBG_P)
        _pptx_rect(slide, 0, 0, 13.33, 0.9, ac)
        _pptx_txt(slide, s_title.upper(), 0.35, 0.1, 12.5, 0.7, bold=True, size=18, color=_C_WHITE_P)
        _pptx_txt(slide, str(i + 1), 12.5, 0.12, 0.6, 0.6, size=11, color=_C_WHITE_P, align=PP_ALIGN.RIGHT)
        _pptx_rect(slide, 0.3, 1.0, 12.73, 5.9, _C_WHITE_P)
        if bullets:
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6))
            txb.word_wrap = True
            tf = txb.text_frame; tf.word_wrap = True
            for j, bullet in enumerate(bullets[:12]):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.space_before = Pt(4)
                dot = p.add_run(); dot.text = "●  "; dot.font.size = Pt(8); dot.font.color.rgb = ac
                run = p.add_run(); run.text = bullet; run.font.size = Pt(12); run.font.color.rgb = _C_TXT_P
        else:
            _pptx_txt(slide, "No data available.", 0.5, 1.2, 12, 0.5, size=11, color=_C_DIM_P)
        _footer(slide)

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()
